"""
Presentation skill: create PowerPoint (PPTX) from outline or natural language.
Manus-style "delivers PowerPoint presentations (PPTX)".
"""
import os
from typing import Dict, Any, List

from viki.skills.base import BaseSkill
from viki.config.logger import viki_logger
from viki.core.utils.path_sandbox import validate_output_path


class PresentationSkill(BaseSkill):
    def __init__(self, controller=None):
        self._controller = controller

    @property
    def name(self) -> str:
        return "presentation"

    @property
    def description(self) -> str:
        return (
            "Create a PowerPoint (.pptx) file. Actions: create_presentation(outline=[{title, bullets}, ...], output_path=...), "
            "or create_from_text(prompt_or_outline=..., output_path=...) to generate from natural language."
        )

    def _create_pptx(self, outline: List[Dict[str, Any]], output_path: str) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        for i, slide_spec in enumerate(outline):
            title = slide_spec.get("title", f"Slide {i + 1}")
            bullets = slide_spec.get("bullets") or slide_spec.get("content") or []
            if isinstance(bullets, str):
                bullets = [bullets]
            layout = prs.slide_layouts[1]  # title and content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            body = slide.placeholders[1].text_frame
            for b in bullets:
                p = body.add_paragraph()
                p.text = b if isinstance(b, str) else str(b)
                p.level = 0
        prs.save(output_path)
        return output_path

    async def execute(self, params: Dict[str, Any]) -> str:
        output_path = params.get("output_path") or params.get("output")
        if not output_path:
            return "Provide output_path= where to save the .pptx file."
        ok, path_or_err = validate_output_path(output_path, controller=self._controller)
        if not ok:
            return path_or_err

        outline = params.get("outline")
        prompt_or_outline = params.get("prompt_or_outline") or params.get("prompt")

        if outline:
            if not isinstance(outline, list):
                return "outline= must be a list of objects with 'title' and 'bullets' (list of strings)."
            try:
                path = self._create_pptx(outline, path_or_err)
                return f"Presentation saved to {path}."
            except Exception as e:
                viki_logger.warning(f"create_presentation: {e}")
                return f"Error creating PPTX: {e}"

        if prompt_or_outline:
            # Use LLM to turn natural language into outline, then create PPTX
            if not self._controller or not hasattr(self._controller, "model_router"):
                return "create_from_text requires controller/model_router. Use create_presentation with outline= instead."
            try:
                model = self._controller.model_router.get_model(capabilities=["general"])
                prompt = (
                    "Convert this into a short PowerPoint outline. Reply with a JSON array of objects, each with "
                    "'title' (string) and 'bullets' (array of strings). No other text. Example: "
                    "[{\"title\": \"Intro\", \"bullets\": [\"Point 1\", \"Point 2\"]}]\n\n"
                    f"User request: {prompt_or_outline}"
                )
                messages = [{"role": "user", "content": prompt}]
                reply = await model.chat(messages, temperature=0.3)
                import json
                # Try to parse JSON from reply (may be wrapped in markdown code block)
                text = reply.strip()
                if "```" in text:
                    start = text.find("[")
                    end = text.rfind("]") + 1
                    if start >= 0 and end > start:
                        text = text[start:end]
                outline_parsed = json.loads(text)
                if not isinstance(outline_parsed, list):
                    outline_parsed = [{"title": "Slide 1", "bullets": [str(outline_parsed)]}]
                path = self._create_pptx(outline_parsed, path_or_err)
                return f"Presentation generated from text and saved to {path}."
            except json.JSONDecodeError as e:
                return f"Could not parse outline from LLM response: {e}"
            except Exception as e:
                viki_logger.warning(f"create_from_text: {e}")
                return f"Error: {e}"

        return "Provide outline= (list of {title, bullets}) or prompt_or_outline= (natural language) and output_path=."
